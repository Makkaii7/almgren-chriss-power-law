"""
Build presentation/slides.pptx — 11-slide deck for the Almgren-Chriss
power-law optimal-execution project.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(REPO_ROOT, "figures")
OUT = os.path.join(REPO_ROOT, "presentation", "slides.pptx")

# ------------------------------------------------------------
# Theme
# ------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x2D, 0x5C)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
DARK = RGBColor(0x22, 0x2B, 0x38)
GRAY = RGBColor(0x55, 0x5E, 0x6B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background() if line is None else None
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=DARK, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_header(slide, title):
    # Navy header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    # Accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), SW, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    # Title text
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))


def add_footer(slide, page_num):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "Optimal Execution — Ali & Mohamed — CODS 612",
             size=10, color=GRAY)
    add_text(slide, Inches(12.3), Inches(7.1), Inches(1), Inches(0.3),
             str(page_num), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_image_fit(slide, path, x, y, w, h):
    if not os.path.exists(path):
        add_text(slide, x, y, w, h, f"[missing figure: {os.path.basename(path)}]",
                 size=14, color=ACCENT)
        return None
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


# ============================================================
# Slide 1 — Title
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.3), SW, Inches(0.08))
stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()

add_text(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.8),
         "Optimal Execution with Non-Linear\n(Power-Law) Market Impact",
         size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

add_text(s, Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.5),
         "CODS 612 — Computational Optimization in Finance",
         size=22, color=RGBColor(0xE7, 0xEC, 0xF5))

add_text(s, Inches(0.8), Inches(5.25), Inches(11.7), Inches(0.5),
         "Authors: Ali & Mohamed",
         size=18, color=RGBColor(0xE7, 0xEC, 0xF5))

add_text(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.5),
         "Professor: Professor Yerkin",
         size=18, color=RGBColor(0xE7, 0xEC, 0xF5))

add_text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.4),
         "Linear Almgren-Chriss baseline vs. empirical power-law impact",
         size=14, color=RGBColor(0xB7, 0xC4, 0xDC))

# ============================================================
# Slide 2 — The Problem
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "The Problem")

add_bullets(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.0), [
    "A fund must liquidate 1,000,000 shares over 50 half-hour periods.",
    "Two competing costs:  sell fast → crash the price (impact);  hold long → exposed to risk (variance).",
    "Goal: find the trading schedule {v_k} that minimizes total expected cost.",
], size=18)

# Objective box
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.3))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT
box.line.color.rgb = NAVY
add_text(s, Inches(1.0), Inches(3.45), Inches(11.3), Inches(0.45),
         "General objective (power-law impact model):",
         size=16, bold=True, color=NAVY)
add_text(s, Inches(1.0), Inches(3.9), Inches(11.3), Inches(0.65),
         "min   Σₖ [  η · |vₖ|^(1+β)   +   γ · σ² · xₖ²  ]",
         size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER, font="Cambria Math")

# Variable legend
add_bullets(s, Inches(0.5), Inches(4.9), Inches(6.2), Inches(2.0), [
    "xₖ   = inventory remaining at period k",
    "vₖ   = xₖ₋₁ − xₖ = shares traded in period k",
    "β    = impact exponent (β = 1 linear, β ≈ 0.6 empirical)",
], size=16)
add_bullets(s, Inches(7.0), Inches(4.9), Inches(5.8), Inches(2.0), [
    "η    = temporary impact coefficient",
    "γ    = risk aversion parameter",
    "σ    = daily price volatility",
], size=16)

add_footer(s, 2)

# ============================================================
# Slide 3 — Linear Baseline (Ali's half)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "The Linear Baseline  (Ali's half)")

add_bullets(s, Inches(0.5), Inches(1.15), Inches(6.3), Inches(4.0), [
    "β = 1.0  ⇒  Quadratic Program:",
    "   min Σ [ η·vₖ² + γ·σ²·xₖ² ]",
    "Closed-form via Euler–Lagrange:",
    "   xₖ = X · sinh(κ(N−k)) / sinh(κN)",
    "κ = √(γ·σ² / η)  controls trajectory shape.",
    "With professor's params the optimum is nearly TWAP (~20,000 shares / period).",
], size=16)

add_image_fit(s, os.path.join(FIG, "linear_trajectory.png"),
              Inches(6.9), Inches(1.2), Inches(6.1), Inches(5.3))
add_footer(s, 3)

# ============================================================
# Slide 4 — Risk Aversion Controls Urgency
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Risk Aversion Controls Urgency")

add_bullets(s, Inches(0.5), Inches(1.15), Inches(4.3), Inches(4.0), [
    "γ is the dial that balances impact cost vs. risk penalty.",
    "Low γ  →  patient, uniform selling (TWAP).",
    "High γ  →  panicked, front-loaded liquidation.",
    "Trade rate shape depends on κ = √(γσ²/η).",
], size=15)

add_image_fit(s, os.path.join(FIG, "linear_trade_rate_sensitivity.png"),
              Inches(4.9), Inches(1.1), Inches(8.1), Inches(5.5))
add_footer(s, 4)

# ============================================================
# Slide 5 — Why Linear Impact is Wrong (Mohamed's half)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Why Linear Impact Is Wrong  (Mohamed's half)")

add_bullets(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.5), [
    "Empirical research (the \"Square Root Law\") shows real market impact follows a power law with β ≈ 0.5 – 0.6.",
    "Linear model (v²) systematically OVERESTIMATES the cost of large trades.",
    "Example:  doubling a trade  —  under linear impact, cost × 4;  under power-law (β = 0.6), cost × 2^1.6 ≈ 3.",
    "The linear model is unnecessarily scared of big trades  ⇒  it over-smooths the schedule.",
], size=18)

# Quote box
qb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(1.5), Inches(5.1), Inches(10.3), Inches(1.4))
qb.fill.solid(); qb.fill.fore_color.rgb = LIGHT
qb.line.color.rgb = ACCENT
add_text(s, Inches(1.8), Inches(5.25), Inches(9.7), Inches(0.5),
         "Square Root Law",
         size=14, bold=True, color=ACCENT)
add_text(s, Inches(1.8), Inches(5.6), Inches(9.7), Inches(0.9),
         "Impact per share ∝ |v|^β with β ≈ 0.5–0.6 across equities, FX, futures.\n"
         "→ Concave impact: big trades are cheaper than linear models predict.",
         size=14, color=DARK)
add_footer(s, 5)

# ============================================================
# Slide 6 — The Power-Law Model
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "The Power-Law Model")

add_bullets(s, Inches(0.5), Inches(1.15), Inches(6.3), Inches(4.0), [
    "With β = 0.6:   min Σ [ η·|vₖ|^1.6 + γ·σ²·xₖ² ]",
    "No closed form — solved numerically.",
    "scipy.optimize.minimize(method='SLSQP').",
    "Decision vars: N−1 intermediate xₖ; x₀ = X, x_N = 0 fixed.",
    "|v| absolute value needed: fractional exponent on negatives is undefined.",
    "Sanity check: β = 1.0 recovers linear solution (max diff < 0.08%).",
], size=15)

add_image_fit(s, os.path.join(FIG, "power_law_trajectory.png"),
              Inches(6.9), Inches(1.2), Inches(6.1), Inches(5.3))
add_footer(s, 6)

# ============================================================
# Slide 7 — Trajectory Comparison
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Trajectory Comparison")

add_image_fit(s, os.path.join(FIG, "trajectory_comparison.png"),
              Inches(0.4), Inches(1.1), Inches(8.4), Inches(5.6))

add_bullets(s, Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.0), [
    "Both solved at professor's baseline parameters.",
    "At β = 0.6 with low γ, curves nearly overlap.",
    "Power-law is slightly more front-loaded.",
    "Baseline cost gap: only $92 (0.10% over optimum).",
], size=14)
add_footer(s, 7)

# ============================================================
# Slide 8 — THE PUNCHLINE: Cost Gap vs. Beta
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "The Punchline:  Cost Gap vs.  β")

add_image_fit(s, os.path.join(FIG, "cost_gap_vs_beta.png"),
              Inches(0.4), Inches(1.1), Inches(8.4), Inches(5.6))

add_bullets(s, Inches(9.0), Inches(1.2), Inches(4.1), Inches(5.5), [
    "As β drops, the gap EXPLODES:",
    "β = 0.6 : $92   (0.1 %)",
    "β = 0.5 : $1,288   (3 %)",
    "β = 0.4 : $3,578   (15.5 %)",
    "β = 0.3 : $10,184   (103 %  — linear costs DOUBLE the optimum!)",
    "Linear model's fear of big trades becomes increasingly wrong as β shrinks.",
], size=13)
add_footer(s, 8)

# ============================================================
# Slide 9 — Gamma Sensitivity
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Gamma Sensitivity:  Linear vs Power-Law")

add_image_fit(s, os.path.join(FIG, "gamma_sensitivity_comparison.png"),
              Inches(0.3), Inches(1.15), Inches(9.5), Inches(5.5))

add_bullets(s, Inches(9.9), Inches(1.3), Inches(3.3), Inches(5.2), [
    "Low γ  →  models agree (both near TWAP).",
    "Medium γ  →  small visible divergence.",
    "High γ  →  models DIVERGE meaningfully.",
    "Model choice matters most when liquidation is urgent.",
], size=13)
add_footer(s, 9)

# ============================================================
# Slide 10 — Conclusion
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Conclusion")

add_bullets(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.0), [
    "Linear Almgren-Chriss is convenient  (QP, closed form)  —  but assumes v² impact.",
    "Real markets follow power-law impact with β ≈ 0.5 – 0.6.",
    "At baseline parameters the linear approximation error is small  (< 0.1%).",
    "But error grows dramatically as:   (1) β decreases,   or   (2) γ increases.",
    "Practical takeaway:  execution desks should calibrate impact empirically, not assume linearity.",
    "Tools used:  Python, NumPy, SciPy (SLSQP), Matplotlib.",
], size=18)
add_footer(s, 10)

# ============================================================
# Slide 11 — References
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "References")

refs = [
    "Almgren, R. and Chriss, N. (2000). Optimal execution of portfolio transactions. Journal of Risk, 3(2), 5–39.",
    "Almgren, R., Thum, C., Hauptmann, E., and Li, H. (2005). Direct estimation of equity market impact. Risk, 18(7), 58–62.",
    "Bouchaud, J.-P. (2010). Price impact. Encyclopedia of Quantitative Finance.",
    "Cornuéjols, G. and Tütüncü, R. (2007). Optimization Methods in Finance. Cambridge University Press.  (course textbook)",
]
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.0))
tf = tb.text_frame; tf.word_wrap = True
for i, r in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    run = p.add_run()
    run.text = f"[{i+1}]   {r}"
    run.font.name = "Calibri"
    run.font.size = Pt(16)
    run.font.color.rgb = DARK

add_footer(s, 11)

# ------------------------------------------------------------
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slide count: {len(prs.slides)}")
